from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE_MD = ROOT / "SYSTEM_USER_GUIDE_PPT_SCRIPT.md"
OUTPUT_PPTX = ROOT / "培训考试系统使用说明_PPT版.pptx"
OUTPUT_PPTX_WITH_IMAGES = ROOT / "培训考试系统使用说明_PPT版_配图.pptx"

IMAGE_POOL = [
    ROOT / "source/vue/xzs-admin/src/assets/logo.png",
    ROOT / "source/vue/xzs-student/src/assets/logo2.png",
    ROOT / "source/vue/xzs-student/src/assets/carousel/1.png",
    ROOT / "source/vue/xzs-student/src/assets/carousel/2.png",
    ROOT / "source/vue/xzs-student/src/assets/carousel/3.png",
    ROOT / "source/vue/xzs-student/src/assets/carousel/4.png",
    ROOT / "source/vue/xzs-student/src/assets/exam-paper/show1.png",
    ROOT / "source/vue/xzs-student/src/assets/exam-paper/show2.png",
    ROOT / "source/vue/xzs-student/src/assets/exam-paper/show3.png",
]

REAL_SCREEN_MAP = {
    "第3页": ROOT / "docs/images/ppt/student_home.png",
    "第5页": ROOT / "docs/images/ppt/admin_dashboard.png",
    "第6页": ROOT / "docs/images/ppt/admin_question_list.png",
    "第8页": ROOT / "docs/images/ppt/admin_question_bank.png",
    "第10页": ROOT / "docs/images/ppt/student_home.png",
    "第11页": ROOT / "docs/images/ppt/student_wrongbook.png",
    "第12页": ROOT / "docs/images/ppt/admin_answer_analysis.png",
    "第13页": ROOT / "docs/images/ppt/admin_answer_analysis.png",
    "第16页": ROOT / "docs/images/ppt/student_record.png",
}


def parse_sections(text: str):
    sections = []
    current_title = None
    current_lines = []

    for line in text.splitlines():
        if line.startswith("## 第"):
            if current_title:
                sections.append((current_title, current_lines))
            current_title = line.replace("## ", "").strip()
            current_lines = []
        else:
            if current_title:
                current_lines.append(line)

    if current_title:
        sections.append((current_title, current_lines))
    return sections


def extract_points(lines):
    points = []
    script = []
    mode = None
    for raw in lines:
        line = raw.strip()
        if not line:
            continue
        if line.startswith("### 页面要点"):
            mode = "points"
            continue
        if line.startswith("### 图示建议"):
            mode = "image"
            continue
        if line.startswith("### 讲解词"):
            mode = "script"
            continue
        if line.startswith("### "):
            mode = None
            continue

        if mode == "points":
            if line.startswith("-") or line[0].isdigit():
                points.append(line.lstrip("- ").strip())
        elif mode == "image":
            if line.startswith("-"):
                points.append(f"【图示】{line.lstrip('- ').strip()}")
        elif mode == "script":
            script.append(line)

    return points[:10], "\n".join(script)


def add_slide(prs: Presentation, title: str, points, notes: str, image_path: Path | None = None):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    if not points:
        points = ["（本页内容请参考手册补充）"]

    for idx, p in enumerate(points):
        para = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        para.text = p
        para.level = 0
        para.font.size = Pt(20)

    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(8.1), Inches(1.4), width=Inches(5.0), height=Inches(4.9))

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def main():
    text = SOURCE_MD.read_text(encoding="utf-8")
    sections = parse_sections(text)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Cover
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "培训考试系统使用说明"
    cover.placeholders[1].text = "PPT自动生成版\n（基于 SYSTEM_USER_GUIDE_PPT_SCRIPT.md）"

    valid_images = [p for p in IMAGE_POOL if p.exists()]

    for idx, (title, lines) in enumerate(sections):
        points, notes = extract_points(lines)
        image = None
        for k, v in REAL_SCREEN_MAP.items():
            if title.startswith(k) and v.exists():
                image = v
                break
        if image is None:
            image = valid_images[idx % len(valid_images)] if valid_images else None
        add_slide(prs, title, points, notes, image)

    prs.save(str(OUTPUT_PPTX))
    prs.save(str(OUTPUT_PPTX_WITH_IMAGES))
    print(f"Generated: {OUTPUT_PPTX}")
    print(f"Generated: {OUTPUT_PPTX_WITH_IMAGES}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
